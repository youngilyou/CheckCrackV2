from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("CheckCrack_AI_Training_Physical_Unit_Metadata_Guide_20260817.pptx")


TITLE = "CheckCrack AI 학습\n원본 픽셀 정량값 및 물리 단위 환산 기록/촬영 가이드"
SUBTITLE = "AI 학습 데이터 생성, 라벨 metadata, 학습 run 기록 범위\n2026-08-17"


SLIDES = [
    (
        "핵심 원칙",
        [
            "AI 학습 단계에서는 실제 mm/cm/m2 단위 결함 치수를 확정 측정하지 않음",
            "학습 단계 정량값은 원본 이미지 픽셀 좌표계를 기준으로 계산",
            "Zoom 화면이나 screenshot은 학습 원본으로 사용하지 않음",
            "실제 crop은 반드시 원본 이미지 파일에서 수행",
            "보고/결과에서 물리 단위가 아니며 pixel 기준임을 명시",
        ],
    ),
    (
        "학습 단계에서 저장할 픽셀 정량값",
        [
            "measurement_basis = original_pixel_coordinate",
            "physical_unit_available = false",
            "bbox_width_px / bbox_height_px",
            "area_px",
            "skeleton_length_px",
            "mean_width_px / max_width_px",
            "polygon_original / polygon_local_px",
            "crop_origin_x/y, crop_width/height, crop_margin",
        ],
    ),
    (
        "Global / Local 좌표 보존 정책",
        [
            "원본 global 좌표 저장 방식은 유지",
            "crop/local mapping metadata를 추가 저장",
            "Global coordinate: Ground Truth 기준 좌표",
            "Local coordinate: crop 내부 YOLO label 생성 기준",
            "둘 중 하나가 아니라 둘 다 보존해야 재현성과 역추적이 가능",
        ],
    ),
    (
        "Crop Metadata 필수 항목",
        [
            "source_image, source_annotation, region_id, defect_type",
            "split, crop_image, label_file",
            "original_width/height",
            "polygon_original",
            "crop_origin_x/y, crop_width/height",
            "polygon_local_px, polygon_yolo",
            "pixel 기준 area/length/width 정량값",
        ],
    ),
    (
        "물리 단위 환산을 위해 반드시 기록할 것",
        [
            "카메라 정보: model, focal length, sensor size, FOV, distortion",
            "촬영 정보: capture_distance_m, 거리 측정 방법, gimbal pitch/yaw",
            "기준 스케일: marker/tile/window 등 실제 길이를 아는 대상",
            "평면 보정: homography, rectification, rectification error",
            "원본 이미지 보존 여부와 EXIF 보존 여부",
        ],
    ),
    (
        "실제 카메라 촬영 설정 가이드",
        [
            "원본 JPEG와 EXIF metadata 보존",
            "5 m / 10 m 등 목표 거리별 촬영 세션 분리",
            "가능하면 외벽면에 정면 촬영",
            "거리 측정 방법을 함께 기록",
            "기준 마커 또는 실제 규격을 아는 구조물 기록",
            "초점, motion blur, 과노출/저노출, 그림자 조건 관리",
        ],
    ),
    (
        "권장 촬영 세션 Metadata",
        [
            "building_id / building_side / capture_session",
            "capture_date / capture_distance_m",
            "distance_measurement_method",
            "camera_model",
            "image_original_preserved",
            "reference_scale_available / reference_object_type",
            "notes: 학습 metadata이며 물리 환산은 조건 충족 시에만 수행",
        ],
    ),
    (
        "Dataset 생성 보완 항목",
        [
            "학습 시작 시 dataset clean rebuild",
            "train == val 구조 제거",
            "deterministic train/val split 적용",
            "향후 building/session split으로 확장 가능하게 설계",
            "metadata.jsonl 생성",
            "split_summary.json 생성",
        ],
    ),
    (
        "GPU / CPU 실행 및 Run Metadata",
        [
            "torch.cuda.is_available() 확인",
            "GPU 있으면 device=0, 없으면 device=cpu",
            "실제 선택 device를 summary/environment에 기록",
            "python, ultralytics, torch, CUDA, GPU 이름 기록",
            "checkpoint sha256, epochs, imgsz, batch 기록",
            "dataset train/val count와 pixel measurement basis 기록",
        ],
    ),
    (
        "픽셀 기반 길이/폭/면적 계산 설계",
        [
            "polygon -> binary mask 생성",
            "area_px = mask pixel count",
            "skeleton_length_px = skeleton 중심선 길이",
            "mean/max_width_px = distance transform 기반 폭 추정",
            "물리 단위 환산은 수행하지 않음",
            "값의 의미: 라벨 품질 및 학습 metadata용 pixel 기준 정량값",
        ],
    ),
    (
        "명시적으로 제외하는 항목",
        [
            "관리 기능",
            "스티칭 파이프라인",
            "분석 화면",
            "보고서 생성 화면",
            "현장 결과 표시",
            "모델 배포 UI",
            "실제 mm/cm 환산 실행",
        ],
    ),
    (
        "최종 구현 방향",
        [
            "원본 global 좌표 저장 유지",
            "crop/local mapping metadata 추가",
            "픽셀 기준 area/length/width 계산",
            "physical unit unavailable 명시",
            "실제 촬영 시 환산 가능 metadata 기록",
            "GPU/CPU 자동 선택",
            "train/val 분리",
            "run별 environment/dataset/model summary 기록",
        ],
    ),
]


def set_run(run, size=24, bold=False, color=(230, 234, 238)):
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.25), Inches(0.65))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run(run, size=28, bold=True, color=(255, 255, 255))


def add_footer(slide, page):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.25), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"CheckCrack AI Training Guide | {page}"
    set_run(run, size=9, color=(130, 144, 155))


def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(Inches(0.85), Inches(1.35), Inches(11.65), Inches(5.35))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(21)
        p.font.color.rgb = RGBColor(224, 231, 235)


def add_accent_bar(slide):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 106, 61)
    shape.line.fill.background()


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # Title slide
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(20, 24, 26)
    add_accent_bar(slide)
    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.55), Inches(11.85), Inches(2.15))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(TITLE.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(34 if i == 0 else 29)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
    sub = slide.shapes.add_textbox(Inches(0.78), Inches(4.1), Inches(10.8), Inches(0.9))
    stf = sub.text_frame
    stf.clear()
    p = stf.paragraphs[0]
    p.text = SUBTITLE
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(170, 183, 190)
    add_footer(slide, 1)

    for idx, (title, bullets) in enumerate(SLIDES, start=2):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(20, 24, 26)
        add_accent_bar(slide)
        add_title(slide, title)
        add_bullets(slide, bullets)
        add_footer(slide, idx)

    prs.save(OUT)
    print(OUT.resolve())


if __name__ == "__main__":
    main()
