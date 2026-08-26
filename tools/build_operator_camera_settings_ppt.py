from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("CheckCrack_Operator_Camera_Capture_Settings_Guide_20260817.pptx")

SLIDES = [
    (
        "현장 촬영 핵심 원칙",
        [
            "원본 JPEG와 EXIF를 보존한다.",
            "촬영 거리와 촬영 세션을 반드시 기록한다.",
            "가능하면 외벽면에 정면으로 촬영한다.",
            "물리 단위 환산이 필요하면 기준 스케일 또는 마커를 함께 촬영한다.",
            "화면 캡처, 리사이즈, 압축 이미지는 학습 원본으로 사용하지 않는다.",
        ],
    ),
    (
        "촬영 전 카메라 설정",
        [
            "촬영 형식: 원본 JPEG 보존",
            "이미지 크기: 최대 해상도",
            "디지털 줌: 사용 금지",
            "초점: 외벽면 기준으로 촬영 전 확인",
            "노출: 과노출과 저노출 방지",
            "셔터속도: 흔들림이 생기지 않도록 충분히 빠르게 설정",
            "ISO: 가능한 낮게 유지",
            "화이트밸런스: 같은 세션에서는 가능하면 고정",
        ],
    ),
    (
        "촬영 거리 설정",
        [
            "5 m 세션과 10 m 세션은 분리한다.",
            "다른 거리는 별도 세션으로 관리한다.",
            "capture_distance_m 값을 기록한다.",
            "거리 측정 방법을 함께 기록한다.",
            "GPS 위치만으로는 외벽면까지의 정확한 거리로 보지 않는다.",
        ],
    ),
    (
        "촬영 자세 및 각도",
        [
            "외벽면에 최대한 정면으로 촬영한다.",
            "같은 외벽면은 유사한 거리와 각도로 촬영한다.",
            "과도한 경사 촬영은 제한한다.",
            "gimbal pitch/yaw, 고도, 외벽면 방향을 기록한다.",
            "경사 촬영 이미지는 물리 환산 신뢰도가 낮을 수 있음을 표시한다.",
        ],
    ),
    (
        "기준 스케일 또는 마커",
        [
            "실제 mm/cm 환산에는 실제 길이를 아는 기준물이 필요하다.",
            "권장: scale marker, ArUco, AprilTag, 균열 게이지",
            "대체: 규격을 아는 타일, 창틀, 외벽 패널, 줄눈 간격",
            "reference_object_type과 reference_length_mm를 기록한다.",
            "기준물이 없으면 픽셀 기준 정량값만 사용한다.",
        ],
    ),
    (
        "이미지 품질 확인",
        [
            "초점이 맞는지 확인한다.",
            "모션 블러가 없는지 확인한다.",
            "과노출과 저노출이 아닌지 확인한다.",
            "균열이 그림자와 겹치지 않는지 확인한다.",
            "외벽 표면이 충분히 선명한지 확인한다.",
            "정상 외벽 이미지도 충분히 포함한다.",
        ],
    ),
    (
        "촬영 세션 이름 규칙",
        [
            "권장 형식: YYYYMMDD_건물ID_동_외벽면_거리",
            "예: 20260817_APT_A_101_EAST_5m",
            "거리가 바뀌면 세션을 분리한다.",
            "외벽면, 건물, 동, 촬영일이 바뀌면 세션을 분리한다.",
            "카메라 설정 또는 기준 마커 유무가 바뀌면 세션을 분리한다.",
        ],
    ),
    (
        "운용자 필수 Metadata",
        [
            "capture_session, capture_date",
            "building_id, building_block, building_side",
            "camera_model, image_original_preserved",
            "capture_distance_m, distance_measurement_method",
            "gimbal_pitch_deg, gimbal_yaw_deg",
            "reference_scale_available, reference_object_type, reference_length_mm",
            "weather, lighting_condition, notes",
        ],
    ),
    (
        "0.3mm 정량측정 검증 전용 준비물 (SmartCrack V2)",
        [
            "폭을 아는 기준 균열 시편(crack gauge) 또는 실측 균열이 있는 외벽 최소 1개소 확보",
            "0.2~0.3 / 0.3~0.5 / 0.5~1.0 / 1.0mm+ 구간별로 가능한 범위에서 각각 참조 균열 확보",
            "RTK 장비 사용 시: RTK FIX 상태(단순 GPS 아님)였는지 세션 metadata에 기록",
            "드론 자체보다 이 기준 시편/참조 균열 확보가 우선 — 기체만 있고 GT가 없으면 0.3mm 성능 입증 불가",
        ],
    ),
    (
        "촬영 후 파일 관리",
        [
            "원본 이미지 폴더를 그대로 보관한다.",
            "리사이즈, 압축, 캡처본과 원본을 섞지 않는다.",
            "세션 metadata 파일을 이미지 폴더와 함께 둔다.",
            "거리별, 외벽면별 폴더를 분리한다.",
            "기준 마커 포함 이미지가 있는지 확인한다.",
            "품질이 낮은 이미지는 삭제하지 말고 warning으로 표시한다.",
        ],
    ),
    (
        "금지 사항",
        [
            "화면 캡처 이미지를 학습 원본으로 사용",
            "메신저 전송으로 압축된 이미지를 학습 원본으로 사용",
            "촬영 거리 미기록",
            "원본 EXIF 삭제",
            "기준 마커 없이 실제 mm/cm 측정값으로 단정",
            "거리, 외벽면, 촬영일이 다른 이미지를 같은 세션으로 혼합",
            "경사 촬영을 정면 촬영과 같은 scale로 취급",
        ],
    ),
    (
        "현장 체크리스트",
        [
            "촬영 전: 원본 JPEG, EXIF, 디지털 줌 미사용, 초점/노출 확인",
            "촬영 전: 촬영 거리 계획, 기준 마커 준비",
            "촬영 중: 5 m / 10 m 세션 분리",
            "촬영 중: 외벽면별 세션 분리, 가능한 정면 촬영",
            "촬영 후: 원본 폴더 보존, metadata 작성",
            "촬영 후: 거리 측정 방법과 품질 warning 기록",
        ],
    ),
    (
        "최종 정리",
        [
            "운용자는 실제 물리 측정을 확정하지 않는다.",
            "운용자는 향후 물리 단위 환산이 가능하도록 촬영 조건과 기준 정보를 남긴다.",
            "AI 학습 단계에서는 원본 이미지 픽셀 기준 정량값만 계산한다.",
            "실제 mm/cm 환산은 원본 이미지, 카메라 정보, 거리, 자세, 기준 스케일이 충분할 때만 가능하다.",
        ],
    ),
]


def style_run(run, size=24, bold=False, color=(230, 234, 238)):
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def add_bar(slide):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 106, 61)
    shape.line.fill.background()


def add_footer(slide, page):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.25), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"CheckCrack Operator Capture Guide | {page}"
    style_run(run, size=9, color=(130, 144, 155))


def add_title(slide, title):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    style_run(run, size=30, bold=True, color=(255, 255, 255))


def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(Inches(0.85), Inches(1.35), Inches(11.65), Inches(5.4))
    tf = box.text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(21)
        p.font.color.rgb = RGBColor(224, 231, 235)
        p.space_after = Pt(8)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(20, 24, 26)
    add_bar(slide)
    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.55), Inches(11.8), Inches(2.1))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "CheckCrack 운용자용\n카메라 촬영 및 기록 설정 가이드"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    sub = slide.shapes.add_textbox(Inches(0.78), Inches(4.1), Inches(10.8), Inches(0.75))
    stf = sub.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sp.text = "AI 학습용 실제 외벽 이미지 촬영 체크리스트 | 2026-08-17"
    sp.font.name = "Malgun Gothic"
    sp.font.size = Pt(19)
    sp.font.color.rgb = RGBColor(170, 183, 190)
    add_footer(slide, 1)

    for idx, (title, bullets) in enumerate(SLIDES, start=2):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(20, 24, 26)
        add_bar(slide)
        add_title(slide, title)
        add_bullets(slide, bullets)
        add_footer(slide, idx)

    prs.save(OUT)
    print(OUT.resolve())


if __name__ == "__main__":
    main()
